#!/usr/bin/env python3
"""Apply supervisor feedback only to Method and Results slides (6–13).

The user's revised source deck is never overwritten. This script writes a new
review copy while preserving all slides outside the requested scope.
"""

from __future__ import annotations

import sys
from pathlib import Path

sys.path.insert(0, "/tmp/mphy_ppt_deps")

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
SOURCE = ROOT / "docs/Pre/Chen_25062145_Slide.pptx"
OUTPUT = ROOT / "docs/Pre/Chen_25062145_Slide_Method_Results_Revised.pptx"
EXAMPLE_DATA = ROOT / "Results/SeventhPhase/EF-001/EF-001_acc_std_vs_raw_room_transitions_detail.png"

FONT = "UCL Sans"

WHITE = "FFFFFF"
OFF_WHITE = "F8F7FA"
PALE_PURPLE = "F0EAF5"
PALE_BLUE = "E8F3FF"
PALE_GREEN = "E7F6EB"
MID_GREY = "667085"
LIGHT_GREY = "D8DCE3"
DARK = "1F2937"
PURPLE = "361A54"
BLUE = "12239E"
RED = "751441"
ORANGE = "E66C37"
GREEN = "1AAB40"


def rgb(value: str) -> RGBColor:
    value = value.lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def set_run_font(run, size: float, colour: str = DARK, bold: bool = False, italic: bool = False) -> None:
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = rgb(colour)


def set_text(
    shape,
    text: str,
    size: float = 20,
    colour: str = DARK,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.MIDDLE,
    margin: float = 0.05,
):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    run = p.add_run()
    run.text = text
    set_run_font(run, size, colour, bold)
    return shape


def add_text(slide, text, x, y, w, h, size=20, colour=DARK, bold=False, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.MIDDLE, margin=0.05):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    return set_text(shape, text, size, colour, bold, align, valign, margin)


def add_box(slide, x, y, w, h, fill=WHITE, line=LIGHT_GREY, line_width=1.1):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(line_width)
    return shape


def add_card(slide, title, body, x, y, w, h, fill, accent):
    add_box(slide, x, y, w, h, fill, accent, 1.2)
    stripe = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(0.09),
        Inches(h),
    )
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = rgb(accent)
    stripe.line.fill.background()
    add_text(slide, title, x + 0.18, y + 0.09, w - 0.30, 0.34, 16, accent, True)
    add_text(slide, body, x + 0.18, y + 0.43, w - 0.30, h - 0.50, 13, DARK, False, valign=MSO_ANCHOR.TOP)


def remove_shapes_except(slide, keep_indices: set[int]) -> None:
    for index, shape in reversed(list(enumerate(slide.shapes))):
        if index in keep_indices:
            continue
        element = shape._element
        element.getparent().remove(element)


def set_alt_text(shape, title: str, description: str) -> None:
    try:
        props = shape._element.xpath(".//p:cNvPr")
        if props:
            props[0].set("name", title)
            props[0].set("descr", description)
    except Exception:
        pass


def revise_method(prs: Presentation) -> None:
    # Slide 6 — state ownership and tools explicitly.
    slide = prs.slides[5]
    set_text(slide.shapes[0], "I built the unified Python pipeline", 32, PURPLE, True, margin=0)
    set_text(
        slide.shapes[2],
        "I devised this evidence hierarchy and implemented it as one reproducible analysis pipeline.",
        17,
        PURPLE,
        True,
        PP_ALIGN.CENTER,
        margin=0.03,
    )
    set_text(slide.shapes[30], "longer when still\nshorter when moving", 14, DARK, False, PP_ALIGN.LEFT)
    set_text(slide.shapes[42], "Implemented and checked by me", 12.5, PURPLE, True)
    set_text(slide.shapes[44], "Reference labels used after prediction", 12.5, RED, True)
    tools = add_box(slide, 8.60, 1.19, 3.72, 0.72, PALE_PURPLE, PURPLE, 1.0)
    set_text(
        tools,
        "Tools: Python · pandas · scikit-learn\nCodex-assisted coding and debugging",
        13.5,
        PURPLE,
        True,
        PP_ALIGN.CENTER,
        margin=0.04,
    )

    # Slide 7 — replace the abstract toy diagram with real raw and processed data.
    slide = prs.slides[6]
    remove_shapes_except(slide, {0, 1})
    set_text(slide.shapes[0], "Worked example: raw sensors → my room estimate", 30, PURPLE, True, margin=0)

    with Image.open(EXAMPLE_DATA) as image:
        iw, ih = image.size
    image_w = 8.55
    image_h = image_w * ih / iw
    picture = slide.shapes.add_picture(
        str(EXAMPLE_DATA),
        Inches(0.55),
        Inches(1.50),
        width=Inches(image_w),
        height=Inches(image_h),
    )
    picture.crop_top = 0.095
    set_alt_text(
        picture,
        "Anonymised raw-motion and room-estimate example",
        "Accelerometer variability above a traditional strongest-RSSI room sequence and the movement-supported room sequence produced by the student's method.",
    )
    cover = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.55),
        Inches(1.50),
        Inches(8.55),
        Inches(0.22),
    )
    cover.fill.solid()
    cover.fill.fore_color.rgb = rgb(WHITE)
    cover.line.fill.background()
    add_text(slide, "Anonymised example: one low-motion episode", 0.78, 1.51, 8.06, 0.28, 12.5, MID_GREY, True, PP_ALIGN.CENTER)

    add_card(slide, "1 · Raw sensor input", "Five-minute accelerometer variability measures wrist movement.", 9.38, 1.50, 3.02, 1.24, PALE_BLUE, BLUE)
    add_card(slide, "2 · Strongest-RSSI rule", "Each window takes the room of the strongest BLE beacon.", 9.38, 2.84, 3.02, 1.24, PALE_PURPLE, PURPLE)
    add_card(slide, "3 · My method", "Movement and episode evidence suppress implausible low-motion switches.", 9.38, 4.18, 3.02, 1.24, PALE_GREEN, GREEN)
    band = add_box(slide, 0.72, 5.86, 11.66, 0.56, OFF_WHITE, PURPLE, 1.0)
    set_text(
        band,
        "Low movement + rapid strongest-RSSI switching suggests radio instability rather than true room changes.",
        16.5,
        PURPLE,
        True,
        PP_ALIGN.CENTER,
        margin=0.04,
    )

    # Slide 8 — reinforce ownership and provenance in the final output.
    slide = prs.slides[7]
    set_text(slide.shapes[0], "Separate outputs: room, behaviour and occupancy", 32, PURPLE, True, margin=0)
    set_text(slide.shapes[28], "Every five-minute row stores the correction reason and evidence source.", 16, PURPLE, True, PP_ALIGN.CENTER, margin=0.03)


def revise_results(prs: Presentation) -> None:
    # Slide 9 — ownership of the cross-dataset run.
    slide = prs.slides[8]
    set_text(slide.shapes[0], "My unified pipeline ran across 13 sessions without failure", 32, PURPLE, True, margin=0)

    # Slide 10 — define baseline vs the student's method and anonymise names.
    slide = prs.slides[9]
    set_text(slide.shapes[0], "My method changed only evidence-supported windows", 32, PURPLE, True, margin=0)
    set_text(slide.shapes[4], "Strongest-RSSI baseline → my method", 18, PURPLE, True)
    set_text(slide.shapes[5], "579 of 25,513 windows\n2.27% of all windows", 19, DARK, False, PP_ALIGN.CENTER)
    set_text(slide.shapes[21], "Dataset 5", 13, MID_GREY, False, PP_ALIGN.CENTER)
    set_text(slide.shapes[24], "Dataset 6", 13, MID_GREY, False, PP_ALIGN.CENTER)
    set_text(slide.shapes[25], "Awake gaps were never filled.\nMost missing evidence stayed missing.", 20, RED, True, PP_ALIGN.CENTER, margin=0.10)

    # Slide 11 — make the comparison explicit and anonymise labelled datasets.
    slide = prs.slides[10]
    set_text(slide.shapes[0], "My method improved agreement over strongest RSSI", 30, PURPLE, True, margin=0)
    for shape_index, label in zip((18, 23, 28, 33), ("Dataset 1", "Dataset 2", "Dataset 3", "Dataset 4")):
        set_text(slide.shapes[shape_index], label, 12, DARK, True, PP_ALIGN.CENTER)
    # Reposition and expand legend labels.
    raw_legend = slide.shapes[35]
    raw_legend.left = Inches(1.43)
    raw_legend.width = Inches(2.08)
    set_text(raw_legend, "Traditional strongest RSSI", 11.5, DARK, False)
    corrected_swatch = slide.shapes[36]
    corrected_swatch.left = Inches(3.76)
    corrected_legend = slide.shapes[37]
    corrected_legend.left = Inches(4.10)
    corrected_legend.width = Inches(2.55)
    set_text(corrected_legend, "My evidence-aware pipeline", 11.5, DARK, False)
    set_text(slide.shapes[40], "What drove the gains?", 15.5, PURPLE, True)
    set_text(
        slide.shapes[41],
        "Datasets 1–3: much of the gain came from probable-away periods matching reference Out labels.\n\nDataset 4: stronger room-level sleep correction.",
        13.5,
        DARK,
        False,
        valign=MSO_ANCHOR.TOP,
        margin=0.03,
    )

    # Slide 12 — anonymise the embedded figure and make the three rows explicit.
    slide = prs.slides[11]
    picture = slide.shapes[2]
    picture.crop_top = max(picture.crop_top, 0.055)
    overlay = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.78),
        Inches(1.34),
        Inches(11.55),
        Inches(0.34),
    )
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = rgb(WHITE)
    overlay.line.fill.background()
    add_text(slide, "Anonymised Dataset 5: main-sleep episodes and RSSI-gap context", 1.25, 1.37, 10.60, 0.28, 13, DARK, True, PP_ALIGN.CENTER)
    set_text(slide.shapes[5], "short supported gaps filled", 12, DARK, False, PP_ALIGN.CENTER)
    set_text(slide.shapes[8], "long unsupported gap retained", 12, DARK, False, PP_ALIGN.CENTER)
    set_text(slide.shapes[11], "strongest RSSI → my method", 12, DARK, False, PP_ALIGN.CENTER)
    set_text(
        slide.shapes[12],
        "Top: strongest-RSSI room\nMiddle: detected sleep\nBottom: gap support",
        13.5,
        PURPLE,
        True,
        PP_ALIGN.CENTER,
        margin=0.03,
    )

    # Slide 13 — anonymise optional-branch datasets and emphasise ownership.
    slide = prs.slides[12]
    set_text(slide.shapes[0], "My optional branches extend location into behaviour", 32, PURPLE, True, margin=0)
    set_text(slide.shapes[23], "Single-floor Dataset 7 correctly selected the K=1 null", 13, RED, True, PP_ALIGN.CENTER)


def validate_scope(source: Presentation, revised: Presentation) -> None:
    """Confirm slides outside 6–13 retain their visible text and geometry."""
    untouched = (0, 1, 2, 3, 4, 13, 14, 15, 16)
    for slide_index in untouched:
        source_slide = source.slides[slide_index]
        revised_slide = revised.slides[slide_index]
        source_shapes = [
            (shape.shape_type, shape.left, shape.top, shape.width, shape.height, getattr(shape, "text", ""))
            for shape in source_slide.shapes
        ]
        revised_shapes = [
            (shape.shape_type, shape.left, shape.top, shape.width, shape.height, getattr(shape, "text", ""))
            for shape in revised_slide.shapes
        ]
        if source_shapes != revised_shapes:
            raise RuntimeError(f"Out-of-scope slide {slide_index + 1} changed")


def main() -> None:
    if not SOURCE.exists():
        raise FileNotFoundError(SOURCE)
    if not EXAMPLE_DATA.exists():
        raise FileNotFoundError(EXAMPLE_DATA)

    source = Presentation(str(SOURCE))
    revised = Presentation(str(SOURCE))
    revise_method(revised)
    revise_results(revised)
    if len(revised.slides) != len(source.slides):
        raise RuntimeError("Slide count changed unexpectedly")
    validate_scope(source, revised)
    revised.save(str(OUTPUT))

    reopened = Presentation(str(OUTPUT))
    validate_scope(source, reopened)
    print(f"Created {OUTPUT}")
    print(f"Slides: {len(reopened.slides)}")
    print("Modified slides: 6–13 only")


if __name__ == "__main__":
    main()
