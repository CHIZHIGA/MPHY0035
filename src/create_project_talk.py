#!/usr/bin/env python3
"""Create the UCL MSc/MRes project-talk deck from the official template.

This script intentionally uses native PowerPoint shapes for most diagrams and
charts so that text remains editable and can request UCL Sans. The supplied UCL
template is preserved as the source of slide masters, layouts, colours and
logos.
"""

from __future__ import annotations

import sys
from pathlib import Path

sys.path.insert(0, "/tmp/mphy_ppt_deps")

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
TEMPLATE = ROOT / "docs/Pre/UCL-PowerPoint-Template-Aptos-2.1.pptx"
OUTPUT = ROOT / "docs/Pre/Chen_25062145_Slide.pptx"
EF002_GAP = ROOT / "Results/SeventhPhase/EF-002/EF-002_sleep_episode_rssi_gap_support.png"

FONT = "UCL Sans"

WHITE = "FFFFFF"
OFF_WHITE = "F8F7FA"
PALE_PURPLE = "F0EAF5"
PALE_BLUE = "E8F3FF"
PALE_GREEN = "E7F6EB"
PALE_ORANGE = "FCEEE8"
MID_GREY = "667085"
LIGHT_GREY = "D8DCE3"
DARK = "1F2937"

PURPLE = "361A54"
BLUE = "12239E"
RED = "751441"
DARK_GREEN = "10433F"
LIGHT_BLUE = "118DFF"
ORANGE = "E66C37"
PINK = "E645AB"
GREEN = "1AAB40"

SLIDE_W = 13.333
SLIDE_H = 7.5


def rgb(value: str) -> RGBColor:
    value = value.lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def remove_all_slides(prs: Presentation) -> None:
    slide_ids = prs.slides._sldIdLst
    for slide_id in list(slide_ids):
        prs.part.drop_rel(slide_id.rId)
        slide_ids.remove(slide_id)


def set_run_font(run, size: float, colour: str = DARK, bold: bool = False, italic: bool = False) -> None:
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = rgb(colour)


def set_text(
    shape,
    text: str,
    size: float = 22,
    colour: str = DARK,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.MIDDLE,
    margin: float = 0.06,
    italic: bool = False,
):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    run = p.add_run()
    run.text = text
    set_run_font(run, size, colour, bold, italic)
    return shape


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    size: float = 22,
    colour: str = DARK,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.MIDDLE,
    margin: float = 0.06,
    italic: bool = False,
):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    return set_text(shape, text, size, colour, bold, align, valign, margin, italic)


def add_rich_text(slide, runs, x, y, w, h, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.MIDDLE, margin=0.05):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    for item in runs:
        run = p.add_run()
        run.text = item[0]
        set_run_font(
            run,
            item[1] if len(item) > 1 else 22,
            item[2] if len(item) > 2 else DARK,
            item[3] if len(item) > 3 else False,
            item[4] if len(item) > 4 else False,
        )
    return shape


def add_box(slide, x, y, w, h, fill=WHITE, line=LIGHT_GREY, radius=True, line_width=1.2):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(line_width)
    return shape


def add_card(slide, title, body, x, y, w, h, fill=WHITE, line=LIGHT_GREY, accent=None, title_size=22, body_size=17):
    box = add_box(slide, x, y, w, h, fill, line)
    if accent:
        stripe = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(x),
            Inches(y),
            Inches(0.10),
            Inches(h),
        )
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = rgb(accent)
        stripe.line.fill.background()
    add_text(slide, title, x + 0.18, y + 0.12, w - 0.32, 0.42, title_size, accent or PURPLE, True)
    add_text(slide, body, x + 0.18, y + 0.57, w - 0.32, h - 0.68, body_size, DARK, False, valign=MSO_ANCHOR.TOP)
    return box


def add_kpi(slide, value, label, x, y, w, h, accent=PURPLE, note=None):
    add_box(slide, x, y, w, h, WHITE, accent, True, 1.5)
    if h < 1.20:
        add_text(slide, value, x + 0.08, y + 0.07, w - 0.16, 0.42, 22, accent, True, PP_ALIGN.CENTER)
        add_text(slide, label, x + 0.10, y + 0.48, w - 0.20, 0.30, 12.5, DARK, False, PP_ALIGN.CENTER)
    else:
        add_text(slide, value, x + 0.08, y + 0.15, w - 0.16, 0.64, 30, accent, True, PP_ALIGN.CENTER)
        add_text(slide, label, x + 0.10, y + 0.82, w - 0.20, 0.50, 17, DARK, False, PP_ALIGN.CENTER)
    if note:
        add_text(slide, note, x + 0.10, y + h - 0.38, w - 0.20, 0.25, 11, MID_GREY, False, PP_ALIGN.CENTER)


def add_circle_label(slide, label, x, y, d, fill, size=17, colour=WHITE, line=None):
    circle = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    circle.fill.solid()
    circle.fill.fore_color.rgb = rgb(fill)
    circle.line.color.rgb = rgb(line or fill)
    set_text(circle, label, size, colour, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE, 0.05)
    return circle


def add_chevron(slide, x, y, w=0.34, h=0.54, colour=MID_GREY):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(colour)
    shape.line.fill.background()
    return shape


def add_line(slide, x1, y1, x2, y2, colour=MID_GREY, width=1.4, dashed=False):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    line.line.color.rgb = rgb(colour)
    line.line.width = Pt(width)
    if dashed:
        line.line.dash_style = 2
    return line


def add_title(slide, title: str, size=34):
    title_shape = slide.shapes.title
    if title_shape is None:
        title_shape = slide.shapes.add_textbox(Inches(0.48), Inches(0.42), Inches(11.1), Inches(0.92))
    set_text(title_shape, title, size, PURPLE, True, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, 0.0)
    return title_shape


def add_source(slide, text: str):
    add_text(slide, text, 0.48, 6.82, 10.55, 0.28, 9.5, MID_GREY, False, valign=MSO_ANCHOR.MIDDLE)


def add_section_tag(slide, text: str, colour=PURPLE):
    box = add_box(slide, 10.82, 0.48, 1.52, 0.34, colour, colour, True, 0)
    font_size = 9.0 if len(text) > 8 else 10.5
    set_text(box, text.upper(), font_size, WHITE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE, 0.01)


def set_alt_text(shape, title: str, description: str) -> None:
    try:
        props = shape._element.xpath(".//p:cNvPr")
        if props:
            props[0].set("name", title)
            props[0].set("descr", description)
    except Exception:
        pass


def add_picture_contain(slide, path: Path, x, y, w, h, alt_title, alt_description):
    with Image.open(path) as image:
        iw, ih = image.size
    image_ratio = iw / ih
    frame_ratio = w / h
    if image_ratio >= frame_ratio:
        draw_w = w
        draw_h = w / image_ratio
        draw_x = x
        draw_y = y + (h - draw_h) / 2
    else:
        draw_h = h
        draw_w = h * image_ratio
        draw_x = x + (w - draw_w) / 2
        draw_y = y
    pic = slide.shapes.add_picture(str(path), Inches(draw_x), Inches(draw_y), Inches(draw_w), Inches(draw_h))
    set_alt_text(pic, alt_title, alt_description)
    return pic


def add_bars(slide, values, x, y, w, h):
    """Native grouped bar chart for raw/corrected agreement values."""
    # Plot grid and labels.
    for tick in (0, 25, 50, 75, 100):
        py = y + h - (tick / 100.0) * h
        add_line(slide, x, py, x + w, py, LIGHT_GREY, 0.8)
        add_text(slide, str(tick), x - 0.46, py - 0.14, 0.38, 0.28, 11, MID_GREY, False, PP_ALIGN.RIGHT)
    add_line(slide, x, y, x, y + h, MID_GREY, 1.1)
    add_line(slide, x, y + h, x + w, y + h, MID_GREY, 1.1)

    group_w = w / len(values)
    bar_w = 0.38
    gap = 0.12
    for index, (label, raw_value, corrected_value) in enumerate(values):
        centre = x + group_w * (index + 0.5)
        for offset, value, colour in (
            (-bar_w - gap / 2, raw_value, PURPLE),
            (gap / 2, corrected_value, GREEN),
        ):
            bh = h * value / 100.0
            bx = centre + offset
            by = y + h - bh
            bar = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                Inches(bx),
                Inches(by),
                Inches(bar_w),
                Inches(bh),
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = rgb(colour)
            bar.line.fill.background()
            add_text(slide, f"{value:.1f}", bx - 0.08, by - 0.34, bar_w + 0.16, 0.28, 11, colour, True, PP_ALIGN.CENTER)
        add_text(slide, label, centre - 0.72, y + h + 0.10, 1.44, 0.55, 12, DARK, True, PP_ALIGN.CENTER)

    # Legend.
    for lx, colour, label in ((x + 0.15, PURPLE, "Raw"), (x + 1.45, GREEN, "Corrected")):
        swatch = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(lx), Inches(y - 0.42), Inches(0.28), Inches(0.22))
        swatch.fill.solid()
        swatch.fill.fore_color.rgb = rgb(colour)
        swatch.line.fill.background()
        add_text(slide, label, lx + 0.34, y - 0.47, 0.92, 0.32, 12, DARK, False)


def add_house(slide, x, y, w, h, two_floor=False):
    body_y = y + 0.55
    body_h = h - 0.55
    body = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(body_y), Inches(w), Inches(body_h))
    body.fill.solid()
    body.fill.fore_color.rgb = rgb(PALE_BLUE)
    body.line.color.rgb = rgb(BLUE)
    body.line.width = Pt(1.4)
    roof = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE, Inches(x - 0.10), Inches(y), Inches(w + 0.20), Inches(0.86))
    roof.fill.solid()
    roof.fill.fore_color.rgb = rgb(PURPLE)
    roof.line.color.rgb = rgb(PURPLE)
    if two_floor:
        add_line(slide, x, body_y + body_h / 2, x + w, body_y + body_h / 2, BLUE, 1.2)
        add_text(slide, "Floor 2", x + 0.08, body_y + 0.05, w - 0.16, body_h / 2 - 0.08, 14, PURPLE, True, PP_ALIGN.CENTER)
        add_text(slide, "Floor 1", x + 0.08, body_y + body_h / 2 + 0.02, w - 0.16, body_h / 2 - 0.08, 14, PURPLE, True, PP_ALIGN.CENTER)
    return body


def set_background(slide, colour=OFF_WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(colour)


def build_deck() -> Presentation:
    prs = Presentation(str(TEMPLATE))
    remove_all_slides(prs)
    prs.core_properties.title = "From room signals to lived experience"
    prs.core_properties.subject = "MSc/MRes Project Talk 2026"
    prs.core_properties.author = "Chen (25062145)"
    prs.core_properties.comments = "Generated from the official UCL PowerPoint template. Text requests UCL Sans."

    title_layout = prs.slide_masters[0].slide_layouts[2]
    content_layout = prs.slide_masters[1].slide_layouts[0]

    # Slide 1: title.
    slide = prs.slides.add_slide(title_layout)
    set_background(slide, OFF_WHITE)
    title = slide.shapes.title
    set_text(title, "From room signals\nto lived experience", 38, PURPLE, True, valign=MSO_ANCHOR.MIDDLE, margin=0)
    add_text(
        slide,
        "Combining wearable activity and RSSI-derived location data in the home",
        2.52,
        2.82,
        7.45,
        0.80,
        23,
        DARK,
        False,
        valign=MSO_ANCHOR.TOP,
    )
    for placeholder in slide.placeholders:
        idx = placeholder.placeholder_format.idx
        if idx == 17:
            set_text(placeholder, "Chen · Student ID 25062145", 20, DARK, False, margin=0)
        elif idx == 10:
            set_text(placeholder, "MSc/MRes Project Talk · July 2026", 17, MID_GREY, False, margin=0)
        elif idx == 11:
            set_text(placeholder, "UCL Department of Medical Physics and Biomedical Engineering", 16, PURPLE, True, margin=0)

    # Slide 2: clinical motivation.
    slide = prs.slides.add_slide(content_layout)
    set_background(slide)
    add_title(slide, "Activity counts need context")
    add_section_tag(slide, "Context")
    add_circle_label(slide, "Wearable\nactivity", 0.75, 2.12, 1.62, PURPLE, 18)
    add_text(slide, "2,000", 2.55, 2.15, 1.55, 0.62, 31, PURPLE, True, PP_ALIGN.CENTER)
    add_text(slide, "steps", 2.55, 2.74, 1.55, 0.38, 18, MID_GREY, False, PP_ALIGN.CENTER)
    add_chevron(slide, 4.28, 2.43, 0.42, 0.65, BLUE)
    add_house(slide, 5.15, 1.75, 2.55, 2.45, True)
    add_text(slide, "Bedroom", 5.28, 2.46, 1.02, 0.38, 14, DARK_GREEN, True, PP_ALIGN.CENTER)
    add_text(slide, "Stairs", 6.43, 2.46, 0.95, 0.38, 14, BLUE, True, PP_ALIGN.CENTER)
    add_text(slide, "Kitchen", 5.30, 3.31, 1.10, 0.38, 14, ORANGE, True, PP_ALIGN.CENTER)
    add_text(slide, "Living", 6.40, 3.31, 0.95, 0.38, 14, RED, True, PP_ALIGN.CENTER)
    add_chevron(slide, 8.02, 2.43, 0.42, 0.65, BLUE)
    add_card(
        slide,
        "Interpretable lived experience",
        "Where?  When?  With whom?\nHow mobile within the home?",
        8.78,
        1.91,
        3.45,
        2.18,
        PALE_GREEN,
        GREEN,
        GREEN,
        21,
        18,
    )
    band = add_box(slide, 0.72, 4.72, 11.50, 1.24, PALE_PURPLE, PURPLE, True, 1.2)
    set_text(
        band,
        "The same movement count can imply different function depending on place, time and social context.",
        23,
        PURPLE,
        True,
        PP_ALIGN.CENTER,
        MSO_ANCHOR.MIDDLE,
        0.16,
    )
    add_source(slide, "Context: Shiffman et al. (2008); Kaye et al. (2011); McLeish et al. (2024).")

    # Slide 3: prior work and gap.
    slide = prs.slides.add_slide(content_layout)
    set_background(slide)
    add_title(slide, "BLE localisation is useful—but real homes are difficult")
    add_section_tag(slide, "Background")
    add_card(
        slide,
        "Transparent baseline",
        "The strongest beacon in a time window provides a simple, calibration-light room proposal.",
        0.62,
        1.66,
        3.65,
        2.50,
        PALE_BLUE,
        BLUE,
        BLUE,
        22,
        18,
    )
    add_card(
        slide,
        "Real-world ambiguity",
        "Body position, mixed transition windows, missing RSSI and cross-floor signals make room estimates unstable.",
        4.67,
        1.66,
        3.65,
        2.50,
        PALE_ORANGE,
        ORANGE,
        ORANGE,
        22,
        18,
    )
    add_card(
        slide,
        "Evaluation gap",
        "Observation reduces realism; diaries are coarse; many datasets support plausibility rather than true accuracy.",
        8.72,
        1.66,
        3.65,
        2.50,
        PALE_PURPLE,
        PURPLE,
        PURPLE,
        22,
        18,
    )
    add_rich_text(
        slide,
        [
            ("Research gap: ", 22, PURPLE, True),
            ("how can multiple sensors add context without hiding uncertainty?", 22, DARK, True),
        ],
        1.02,
        4.73,
        11.10,
        0.88,
        PP_ALIGN.CENTER,
    )
    add_source(slide, "BLE context: Bai et al. (2020); Roohi & Fekr (2025); floor sensing: Falcon & Schulzrinne (2017).")

    # Slide 4: aim and objectives.
    slide = prs.slides.add_slide(content_layout)
    set_background(slide)
    add_title(slide, "Aim: move from room labels to behavioural evidence")
    add_section_tag(slide, "Aim")
    aim_box = add_box(slide, 0.74, 1.42, 11.58, 0.88, PALE_PURPLE, PURPLE, True, 1.3)
    set_text(
        aim_box,
        "Combine wearable movement, BLE RSSI and environmental context to describe behaviour in the home.",
        21,
        PURPLE,
        True,
        PP_ALIGN.CENTER,
        margin=0.10,
    )
    objectives = [
        ("1", "Align", "Reproducible\npreprocessing", BLUE),
        ("2", "Baseline", "Transparent\nRSSI location", PURPLE),
        ("3", "Fuse", "Movement +\nfloor context", ORANGE),
        ("4", "Translate", "Behavioural\nsummaries", GREEN),
        ("5", "Evaluate", "Evidence matched\nto each dataset", RED),
    ]
    x_positions = [0.62, 3.11, 5.60, 8.09, 10.58]
    for x, (number, title, body, colour) in zip(x_positions, objectives):
        add_circle_label(slide, number, x + 0.73, 2.67, 0.70, colour, 22)
        add_text(slide, title, x, 3.48, 2.16, 0.42, 19, colour, True, PP_ALIGN.CENTER)
        add_text(slide, body, x, 3.93, 2.16, 0.94, 17, DARK, False, PP_ALIGN.CENTER)
    add_line(slide, 1.72, 3.02, 11.67, 3.02, LIGHT_GREY, 1.8)
    add_text(
        slide,
        "Location is treated as intermediate evidence—not the final clinical outcome.",
        2.15,
        5.42,
        9.02,
        0.74,
        22,
        DARK_GREEN,
        True,
        PP_ALIGN.CENTER,
    )

    # Slide 5: data and evidence levels.
    slide = prs.slides.add_slide(content_layout)
    set_background(slide)
    add_title(slide, "Heterogeneous data require evidence-aware evaluation")
    add_section_tag(slide, "Data")
    inputs = [
        ("RSSI", "required", PURPLE, 15),
        ("ACC /\nsteps", "required", BLUE, 13),
        ("Pressure", "optional", ORANGE, 12),
        ("Paired\nuser", "optional", GREEN, 13),
        ("Reference\nlabels", "optional", RED, 10.5),
    ]
    for i, (name, status, colour, label_size) in enumerate(inputs):
        x = 0.63 + i * 2.48
        add_circle_label(slide, name, x + 0.48, 1.54, 1.34, colour, label_size)
        add_text(slide, status, x + 0.32, 2.90, 1.66, 0.33, 13, MID_GREY, False, PP_ALIGN.CENTER)
    add_text(slide, "Evidence available", 0.66, 3.61, 2.25, 0.42, 18, PURPLE, True)
    tiers = [
        ("No reference", "Transition counts, provenance and visual plausibility", PALE_BLUE, BLUE),
        ("Existing annotations", "Coverage, balanced accuracy, macro-F1 and agreement", PALE_GREEN, GREEN),
        ("Independent validation", "Required after rules are frozen; not yet available", PALE_ORANGE, ORANGE),
    ]
    for i, (title, body, fill, colour) in enumerate(tiers):
        add_card(slide, title, body, 0.67 + i * 4.10, 4.05, 3.72, 1.63, fill, colour, colour, 19, 15)
    add_text(
        slide,
        "Reference labels are loaded only after predictions and model parameters are generated.",
        1.58,
        5.96,
        10.15,
        0.50,
        17,
        RED,
        True,
        PP_ALIGN.CENTER,
    )

    # Slide 6: pipeline.
    slide = prs.slides.add_slide(content_layout)
    set_background(slide)
    add_title(slide, "One auditable five-minute pipeline")
    add_section_tag(slide, "Method")
    add_text(
        slide,
        "RSSI proposes location; contextual evidence determines whether it is accepted, stabilised, constrained or unresolved.",
        0.72,
        1.28,
        11.68,
        0.62,
        18,
        PURPLE,
        True,
        PP_ALIGN.CENTER,
    )
    pipeline_top = [
        ("Raw inputs", "RSSI · movement\noptional context", PALE_BLUE, BLUE),
        ("Capability audit", "Sensors · labels\nmetadata", WHITE, PURPLE),
        ("Five-minute grid", "UTC alignment\nlocal time retained", WHITE, BLUE),
        ("Features", "RSSI strength\nACC SD or steps", PALE_GREEN, GREEN),
    ]
    for i, (title, body, fill, colour) in enumerate(pipeline_top):
        x = 0.55 + i * 3.10
        add_card(slide, title, body, x, 2.05, 2.55, 1.34, fill, colour, colour, 18, 15)
        if i < len(pipeline_top) - 1:
            add_chevron(slide, x + 2.67, 2.47, 0.30, 0.48, MID_GREY)
    pipeline_bottom = [
        ("Movement + state masks", "sleep · occupancy · pressure", PALE_ORANGE, ORANGE),
        ("Adaptive RSSI", "trailing window by movement", PALE_PURPLE, PURPLE),
        ("Final timeline", "room · behaviour · occupancy", PALE_BLUE, BLUE),
        ("Downstream outputs", "co-presence · agreement · mobility", PALE_GREEN, GREEN),
    ]
    for i, (title, body, fill, colour) in enumerate(pipeline_bottom):
        x = 0.55 + i * 3.10
        add_card(slide, title, body, x, 4.25, 2.55, 1.34, fill, colour, colour, 18, 15)
        if i < len(pipeline_bottom) - 1:
            add_chevron(slide, x + 2.67, 4.67, 0.30, 0.48, MID_GREY)
    add_line(slide, 11.72, 3.40, 11.72, 4.22, MID_GREY, 1.6)
    add_text(slide, "Prediction", 0.77, 5.97, 2.15, 0.30, 13, MID_GREY, True)
    add_line(slide, 2.08, 6.12, 8.65, 6.12, LIGHT_GREY, 1.2)
    add_text(slide, "Post-hoc evaluation", 8.82, 5.97, 2.58, 0.30, 13, RED, True)

    # Slide 7: movement-aware windows.
    slide = prs.slides.add_slide(content_layout)
    set_background(slide)
    add_title(slide, "Movement changes how RSSI is trusted")
    add_section_tag(slide, "Method")
    add_card(slide, "Low movement", "Longer trailing history suppresses isolated radio instability.", 0.62, 1.56, 5.86, 1.14, PALE_BLUE, BLUE, BLUE, 23, 17)
    add_card(slide, "Higher movement", "Shorter history preserves plausible room transitions.", 6.84, 1.56, 5.86, 1.14, PALE_ORANGE, ORANGE, ORANGE, 23, 17)
    # Mini timelines.
    left_colours = [PURPLE, PURPLE, PURPLE, PINK, PURPLE, PURPLE]
    right_colours = [PURPLE, PURPLE, GREEN, GREEN, GREEN, GREEN]
    for i, colour in enumerate(left_colours):
        add_circle_label(slide, "", 1.07 + i * 0.68, 3.24, 0.30, colour, 1)
    for i, colour in enumerate(right_colours):
        add_circle_label(slide, "", 7.31 + i * 0.68, 3.24, 0.30, colour, 1)
    add_line(slide, 1.07, 3.88, 4.77, 3.88, BLUE, 3)
    add_text(slide, "30 min trailing window", 1.24, 3.97, 3.38, 0.35, 15, BLUE, True, PP_ALIGN.CENTER)
    add_line(slide, 9.96, 3.88, 10.98, 3.88, ORANGE, 3)
    add_text(slide, "5–15 min", 9.68, 3.97, 1.54, 0.35, 15, ORANGE, True, PP_ALIGN.CENTER)
    add_text(slide, "Transient RSSI winner", 3.18, 2.85, 2.05, 0.32, 13, PINK, True, PP_ALIGN.CENTER)
    add_text(slide, "Room transition retained", 8.78, 2.85, 2.78, 0.32, 13, GREEN, True, PP_ALIGN.CENTER)
    principles = [
        ("Movement is context", "not a direct room classifier", PURPLE),
        ("Awake inference is causal", "current and past evidence only", BLUE),
        ("Thresholds are fitted per session", "using one common clustering rule", GREEN),
    ]
    for i, (title, body, colour) in enumerate(principles):
        add_card(slide, title, body, 0.76 + i * 4.08, 4.82, 3.66, 1.15, WHITE, colour, colour, 17, 14)

    # Slide 8: state hierarchy and safeguards.
    slide = prs.slides.add_slide(content_layout)
    set_background(slide)
    add_title(slide, "Room, behaviour and occupancy remain separate")
    add_section_tag(slide, "Method")
    outputs = [
        ("ROOM", "Raw RSSI room\nCorrected room\nCorrection reason", PALE_BLUE, BLUE),
        ("BEHAVIOUR", "Awake\nMain sleep\nMovement unresolved", PALE_PURPLE, PURPLE),
        ("OCCUPANCY", "Indoor observed\nProbable away\nUnknown", PALE_GREEN, GREEN),
    ]
    for i, (title, body, fill, colour) in enumerate(outputs):
        box = add_box(slide, 0.72 + i * 4.13, 1.48, 3.72, 2.29, fill, colour, True, 1.4)
        add_text(slide, title, 0.90 + i * 4.13, 1.68, 3.36, 0.42, 19, colour, True, PP_ALIGN.CENTER)
        add_text(slide, body, 0.92 + i * 4.13, 2.17, 3.32, 1.28, 17, DARK, False, PP_ALIGN.CENTER)
    add_text(slide, "Safeguards", 0.74, 4.10, 1.54, 0.40, 19, PURPLE, True)
    safeguards = [
        ("Awake gaps", "never filled", BLUE),
        ("Sleep gaps", "need two-sided support", PURPLE),
        ("Pressure", "K=1 null allowed", ORANGE),
        ("Weak evidence", "remains unresolved", GREEN),
    ]
    for i, (title, body, colour) in enumerate(safeguards):
        add_card(slide, title, body, 0.72 + i * 3.04, 4.55, 2.70, 1.10, WHITE, colour, colour, 16, 14)
    band = add_box(slide, 1.18, 5.95, 10.60, 0.48, PALE_PURPLE, PALE_PURPLE, True, 0)
    set_text(band, "The pipeline reports uncertainty instead of repairing it manually.", 17, PURPLE, True, PP_ALIGN.CENTER, margin=0.03)

    # Slide 9: audit results.
    slide = prs.slides.add_slide(content_layout)
    set_background(slide)
    add_title(slide, "The unified pipeline ran across 13 sessions without failure")
    add_section_tag(slide, "Results")
    kpis = [
        ("10", "collections", BLUE),
        ("13", "participant-sessions", PURPLE),
        ("25,513", "five-minute windows", GREEN),
        ("2,126.1 h", "recorded time", ORANGE),
        ("0", "pipeline failures", RED),
    ]
    for i, (value, label, colour) in enumerate(kpis):
        add_kpi(slide, value, label, 0.55 + i * 2.51, 1.62, 2.22, 1.57, colour)
    add_card(slide, "Movement", "9 ACC sessions\n4 step fallbacks", 0.73, 3.78, 3.57, 1.43, PALE_BLUE, BLUE, BLUE, 20, 17)
    add_card(slide, "Main sleep", "9 sessions resolved\n4 left unresolved", 4.87, 3.78, 3.57, 1.43, PALE_PURPLE, PURPLE, PURPLE, 20, 17)
    add_card(slide, "Away modelling", "13 sessions resolved\n158 probable-away runs", 9.01, 3.78, 3.57, 1.43, PALE_GREEN, GREEN, GREEN, 20, 17)
    add_text(
        slide,
        "Excluded input ≠ unresolved output: unresolved states are part of the result.",
        1.14,
        5.72,
        10.75,
        0.57,
        19,
        PURPLE,
        True,
        PP_ALIGN.CENTER,
    )

    # Slide 10: selective correction results.
    slide = prs.slides.add_slide(content_layout)
    set_background(slide)
    add_title(slide, "Corrections were selective—not blanket smoothing")
    add_section_tag(slide, "Results")
    add_card(slide, "Windows changed", "579 of 25,513\n2.27% of all windows", 0.68, 1.54, 3.70, 1.86, PALE_PURPLE, PURPLE, PURPLE, 22, 20)
    # Proportion bar.
    add_box(slide, 0.96, 3.05, 3.14, 0.18, LIGHT_GREY, LIGHT_GREY, False, 0)
    add_box(slide, 0.96, 3.05, 3.14 * 0.0227, 0.18, PURPLE, PURPLE, False, 0)
    add_card(slide, "Observed transitions", "3,353  →  2,961\n−392 transitions (−11.7%)", 4.80, 1.54, 3.70, 1.86, PALE_BLUE, BLUE, BLUE, 22, 20)
    add_card(slide, "Room coverage", "76.48%  →  76.70%\n+0.22 percentage points", 8.92, 1.54, 3.70, 1.86, PALE_GREEN, GREEN, GREEN, 22, 20)
    add_text(slide, "Evidence-supported sleep-gap filling", 0.73, 4.12, 3.62, 0.42, 18, ORANGE, True)
    add_kpi(slide, "57", "five-minute gaps filled", 0.73, 4.60, 2.40, 1.46, ORANGE)
    add_text(slide, "51", 3.50, 4.66, 1.15, 0.48, 27, PURPLE, True, PP_ALIGN.CENTER)
    add_text(slide, "EF-002", 3.50, 5.12, 1.15, 0.32, 14, MID_GREY, False, PP_ALIGN.CENTER)
    add_text(slide, "+", 4.69, 4.78, 0.42, 0.42, 22, MID_GREY, True, PP_ALIGN.CENTER)
    add_text(slide, "6", 5.16, 4.66, 1.15, 0.48, 27, GREEN, True, PP_ALIGN.CENTER)
    add_text(slide, "KM PanH", 5.08, 5.12, 1.32, 0.32, 14, MID_GREY, False, PP_ALIGN.CENTER)
    statement = add_box(slide, 7.03, 4.30, 5.27, 1.49, WHITE, RED, True, 1.2)
    set_text(
        statement,
        "Awake gaps were never filled.\nMost missing evidence stayed missing.",
        21,
        RED,
        True,
        PP_ALIGN.CENTER,
        margin=0.12,
    )

    # Slide 11: labelled agreement.
    slide = prs.slides.add_slide(content_layout)
    set_background(slide)
    add_title(slide, "End-to-end agreement improved in all four labelled datasets")
    add_section_tag(slide, "Results")
    values = [
        ("DH Paris", 28.1, 61.9),
        ("DH PanoH", 56.3, 89.6),
        ("DH Strad", 57.9, 82.0),
        ("KM Mal", 57.1, 85.8),
    ]
    add_bars(slide, values, 0.94, 1.86, 8.62, 3.86)
    add_card(slide, "What drove the gain?", "DH datasets: much of the gain came from representing probable away and matching reference Out periods.\n\nKM Mal: stronger room-level sleep correction.", 9.82, 1.65, 2.76, 3.90, PALE_PURPLE, PURPLE, PURPLE, 19, 15)
    add_text(
        slide,
        "Agreement with existing annotations—not independently verified ground truth.",
        1.15,
        6.36,
        10.55,
        0.36,
        14,
        RED,
        True,
        PP_ALIGN.CENTER,
    )

    # Slide 12: EF-002 case study.
    slide = prs.slides.add_slide(content_layout)
    set_background(slide)
    add_title(slide, "Case study: only supported RSSI gaps were repaired")
    add_section_tag(slide, "Results")
    add_picture_contain(
        slide,
        EF002_GAP,
        0.54,
        1.34,
        12.02,
        4.22,
        "EF-002 sleep-gap evidence timeline",
        "Timeline showing raw room, main-sleep episodes and RSSI gaps. Short bedroom-supported gaps are marked separately from a long unsupported gap.",
    )
    add_kpi(slide, "51", "supported windows filled", 0.76, 5.72, 2.65, 0.92, GREEN)
    add_kpi(slide, "≈700 min", "unsupported gap retained", 3.76, 5.72, 2.80, 0.92, ORANGE)
    add_kpi(slide, "543 → 513", "room transitions", 6.91, 5.72, 2.60, 0.92, BLUE)
    add_text(slide, "Conservative correction is a feature, not a failure.", 9.74, 5.78, 2.58, 0.69, 17, PURPLE, True, PP_ALIGN.CENTER)

    # Slide 13: optional branches.
    slide = prs.slides.add_slide(content_layout)
    set_background(slide)
    add_title(slide, "Optional branches extend location into behaviour")
    add_section_tag(slide, "Results")
    # Left: two-person.
    add_box(slide, 0.64, 1.48, 5.86, 4.54, PALE_GREEN, GREEN, True, 1.3)
    add_text(slide, "Two-person co-presence", 0.91, 1.70, 5.30, 0.47, 22, GREEN, True, PP_ALIGN.CENTER)
    add_circle_label(slide, "P1", 1.32, 2.40, 0.76, PURPLE, 17)
    add_circle_label(slide, "P2", 2.31, 2.40, 0.76, BLUE, 17)
    add_house(slide, 3.76, 2.19, 1.69, 1.54, False)
    add_text(slide, "Independent inference\nthen shared-room comparison", 0.95, 3.37, 2.50, 0.77, 16, DARK, True, PP_ALIGN.CENTER)
    add_text(slide, "32.67 h", 3.42, 3.95, 2.37, 0.55, 29, GREEN, True, PP_ALIGN.CENTER)
    add_text(slide, "corrected same-room time", 3.36, 4.46, 2.50, 0.40, 14, DARK, False, PP_ALIGN.CENTER)
    add_text(slide, "+41 added  ·  −48 removed windows", 1.19, 4.90, 4.74, 0.43, 17, PURPLE, True, PP_ALIGN.CENTER)
    add_text(slide, "Descriptive co-presence—not verified social interaction", 1.04, 5.46, 5.12, 0.34, 13, RED, True, PP_ALIGN.CENTER)
    # Right: pressure.
    add_box(slide, 6.84, 1.48, 5.86, 4.54, PALE_BLUE, BLUE, True, 1.3)
    add_text(slide, "Pressure-supported floor constraint", 7.10, 1.70, 5.36, 0.47, 22, BLUE, True, PP_ALIGN.CENTER)
    add_house(slide, 7.46, 2.18, 2.02, 1.89, True)
    add_circle_label(slide, "wearable", 10.28, 2.48, 1.04, ORANGE, 13)
    add_chevron(slide, 9.66, 2.74, 0.36, 0.55, BLUE)
    add_text(slide, "K = 2", 9.84, 3.68, 1.78, 0.42, 24, BLUE, True, PP_ALIGN.CENTER)
    add_text(slide, "silhouette 0.865", 9.69, 4.08, 2.08, 0.34, 14, DARK, False, PP_ALIGN.CENTER)
    add_text(slide, "5 pressure-only changes · 7 overall", 7.40, 4.74, 4.71, 0.42, 17, PURPLE, True, PP_ALIGN.CENTER)
    add_text(slide, "Single-floor Home_X001 correctly selected the K=1 null", 7.17, 5.31, 5.23, 0.40, 13, RED, True, PP_ALIGN.CENTER)
    add_text(slide, "Localisation becomes intermediate evidence for behavioural summaries.", 1.36, 6.26, 10.62, 0.42, 18, PURPLE, True, PP_ALIGN.CENTER)

    # Slide 14: interpretation and limitations.
    slide = prs.slides.add_slide(content_layout)
    set_background(slide)
    add_title(slide, "Generalisability came with an explicit trade-off")
    add_section_tag(slide, "Discussion")
    add_card(
        slide,
        "What the unified method gains",
        "• One schema across heterogeneous sessions\n• Explicit missing and unresolved states\n• Correction provenance and optional branches\n• Comparable cross-dataset reporting",
        0.66,
        1.48,
        5.72,
        3.78,
        PALE_GREEN,
        GREEN,
        GREEN,
        23,
        18,
    )
    add_card(
        slide,
        "What the evidence cannot yet claim",
        "• Not always the best dataset-specific classifier\n• Limited independent reference labels\n• Step absence is weaker than ACC\n• Multi-floor behaviour tested in one 80-hour case",
        6.75,
        1.48,
        5.72,
        3.78,
        PALE_ORANGE,
        ORANGE,
        ORANGE,
        23,
        18,
    )
    implication = add_box(slide, 1.08, 5.60, 11.02, 0.76, PALE_PURPLE, PURPLE, True, 1.2)
    set_text(
        implication,
        "Useful for exploratory, interpretable home-behaviour analysis—not yet a validated clinical system.",
        20,
        PURPLE,
        True,
        PP_ALIGN.CENTER,
        margin=0.10,
    )

    # Slide 15: take-home.
    slide = prs.slides.add_slide(content_layout)
    set_background(slide)
    add_title(slide, "Take-home message")
    add_section_tag(slide, "Conclusion")
    takeaways = [
        ("1", "Context", "Room-level evidence makes wearable activity more interpretable.", BLUE),
        ("2", "Restraint", "Sensor fusion changes estimates only when supporting evidence is sufficient.", PURPLE),
        ("3", "Translation", "The same timeline supports co-presence, floor and mobility summaries.", GREEN),
    ]
    for i, (number, heading, body, colour) in enumerate(takeaways):
        y = 1.48 + i * 1.37
        add_circle_label(slide, number, 0.78, y + 0.10, 0.72, colour, 22)
        add_text(slide, heading, 1.82, y, 2.10, 0.46, 22, colour, True)
        add_text(slide, body, 4.12, y, 8.06, 0.66, 19, DARK, False)
    add_text(slide, "Next: freeze the rules and test them on independent labelled homes.", 1.22, 5.67, 10.86, 0.52, 19, RED, True, PP_ALIGN.CENTER)
    thanks = add_box(slide, 3.80, 6.31, 5.60, 0.56, PURPLE, PURPLE, True, 0)
    set_text(thanks, "Thank you — questions", 20, WHITE, True, PP_ALIGN.CENTER, margin=0.04)

    # Slide 16: backup thresholds.
    slide = prs.slides.add_slide(content_layout)
    set_background(slide)
    add_title(slide, "Backup: replication-critical thresholds")
    add_section_tag(slide, "Backup", MID_GREY)
    rows = [
        ("Common inference unit", "Fixed non-overlapping 5-minute windows"),
        ("Movement clustering", "K = 2–4; silhouette ≥ 0.25; each cluster ≥ 5% and ≥ 20 windows"),
        ("Main-sleep candidate", "≥ 60 min and ≥ 60% low-motion windows"),
        ("Sleep-room dominance", "Dominant room ≥ 60% of observed episode windows"),
        ("Sleep-gap fill", "Both adjacent 30-min contexts support the same room by ≥ two thirds"),
        ("Pressure grouping", "Silhouette ≥ 0.75; adjacent separation ≥ 0.30 hPa and > 3 pooled MADs"),
        ("Floor constraint", "Pressure confidence ≥ 0.75 plus movement support"),
    ]
    for i, (rule, threshold) in enumerate(rows):
        y = 1.48 + i * 0.70
        fill = WHITE if i % 2 == 0 else PALE_BLUE
        add_box(slide, 0.70, y, 11.78, 0.60, fill, LIGHT_GREY, False, 0.7)
        add_text(slide, rule, 0.88, y + 0.04, 3.22, 0.50, 15, PURPLE, True)
        add_text(slide, threshold, 4.17, y + 0.04, 8.08, 0.50, 15, DARK, False)
    add_text(slide, "All failure checks return unresolved or inactive states; they are not manually overridden.", 1.26, 6.50, 10.58, 0.36, 14, RED, True, PP_ALIGN.CENTER)

    # Slide 17: selected sources.
    slide = prs.slides.add_slide(content_layout)
    set_background(slide)
    add_title(slide, "Backup: selected sources")
    add_section_tag(slide, "Backup", MID_GREY)
    refs_left = (
        "[1] Shiffman, Stone & Hufford. Ecological momentary assessment. Annual Review of Clinical Psychology, 2008.\n\n"
        "[2] Kaye et al. Home-based, unobtrusive and continuous assessment of ageing. Journals of Gerontology B, 2011.\n\n"
        "[3] Sadeh & Acebo. The role of actigraphy in sleep medicine. Sleep Medicine Reviews, 2002.\n\n"
        "[4] McLeish et al. Validating room location within the home for functional measures from digital health technologies. IEEE EMBC, 2024."
    )
    refs_right = (
        "[5] Bai et al. A low-cost indoor positioning system using Bluetooth Low Energy. IEEE Access, 2020.\n\n"
        "[6] Roohi & Fekr. A comparative analysis of indoor localisation technologies. Computer Networks, 2025.\n\n"
        "[7] Falcon & Schulzrinne. Predicting floor level using smartphone sensor data. arXiv:1710.11122, 2017.\n\n"
        "Project evidence: FinalReport and EighthPhase cross-dataset outputs in this repository."
    )
    add_text(slide, refs_left, 0.68, 1.45, 5.78, 4.97, 15, DARK, False, valign=MSO_ANCHOR.TOP)
    add_text(slide, refs_right, 6.82, 1.45, 5.68, 4.97, 15, DARK, False, valign=MSO_ANCHOR.TOP)

    return prs


def validate(prs: Presentation) -> None:
    if len(prs.slides) != 17:
        raise RuntimeError(f"Expected 17 slides, found {len(prs.slides)}")
    for slide_number, slide in enumerate(prs.slides, 1):
        if slide_number != 1 and slide.shapes.title is None:
            raise RuntimeError(f"Slide {slide_number} has no title")
        for shape in slide.shapes:
            if shape.left < 0 or shape.top < 0:
                raise RuntimeError(f"Slide {slide_number}: negative shape position: {shape.name}")
            if shape.left + shape.width > prs.slide_width + Inches(0.01):
                raise RuntimeError(f"Slide {slide_number}: shape exceeds width: {shape.name}")
            if shape.top + shape.height > prs.slide_height + Inches(0.01):
                raise RuntimeError(f"Slide {slide_number}: shape exceeds height: {shape.name}")


def main() -> None:
    if not TEMPLATE.exists():
        raise FileNotFoundError(TEMPLATE)
    if not EF002_GAP.exists():
        raise FileNotFoundError(EF002_GAP)
    prs = build_deck()
    validate(prs)
    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    # Re-open to verify the OPC package and relationships are readable.
    check = Presentation(str(OUTPUT))
    validate(check)
    print(f"Created {OUTPUT}")
    print(f"Slides: {len(check.slides)}")


if __name__ == "__main__":
    main()
